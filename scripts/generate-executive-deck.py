#!/usr/bin/env python3
"""Generate Stillform applied-science executive deck."""

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(10, 13, 19)
SURFACE = RGBColor(19, 25, 36)
SURFACE_ALT = RGBColor(27, 34, 47)
LINE = RGBColor(59, 69, 87)
TEXT = RGBColor(236, 240, 245)
TEXT_DIM = RGBColor(167, 178, 194)
ACCENT = RGBColor(235, 171, 74)
BLUE = RGBColor(98, 159, 246)
GREEN = RGBColor(101, 191, 149)
RED = RGBColor(216, 97, 97)
PURPLE = RGBColor(164, 125, 235)

FONT = "Calibri"
MONO = "Consolas"


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, idx, total):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.38))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SURFACE
    bar.line.fill.background()

    k = slide.shapes.add_textbox(Inches(0.56), Inches(0.12), Inches(8.5), Inches(0.2))
    kp = k.text_frame.paragraphs[0]
    kr = kp.add_run()
    kr.text = "STILLFORM · APPLIED SCIENCE DECK"
    kr.font.name = MONO
    kr.font.size = Pt(9)
    kr.font.color.rgb = ACCENT

    n = slide.shapes.add_textbox(Inches(11.2), Inches(0.12), Inches(1.6), Inches(0.2))
    np = n.text_frame.paragraphs[0]
    np.alignment = PP_ALIGN.RIGHT
    nr = np.add_run()
    nr.text = f"{idx:02d} / {total:02d}"
    nr.font.name = MONO
    nr.font.size = Pt(9)
    nr.font.color.rgb = TEXT_DIM


def add_title(slide, title, subtitle):
    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.72), Inches(12.0), Inches(1.2))
    tp = t.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.name = FONT
    tr.font.size = Pt(40)
    tr.font.bold = True
    tr.font.color.rgb = TEXT

    s = slide.shapes.add_textbox(Inches(0.74), Inches(1.62), Inches(11.95), Inches(0.45))
    sp = s.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = subtitle
    sr.font.name = FONT
    sr.font.size = Pt(16)
    sr.font.color.rgb = TEXT_DIM


def add_source(slide, text):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(7.07), Inches(12.0), Inches(0.24))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"Source: {text}"
    r.font.name = MONO
    r.font.size = Pt(9)
    r.font.color.rgb = TEXT_DIM


def add_card(slide, x, y, w, h, alt=False):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE_ALT if alt else SURFACE
    card.line.color.rgb = LINE
    card.line.width = Pt(1.1)
    return card


def style_chart(chart, legend=True):
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = FONT
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = TEXT_DIM

    if chart.category_axis:
        chart.category_axis.tick_labels.font.name = FONT
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.category_axis.tick_labels.font.color.rgb = TEXT_DIM
    if chart.value_axis:
        chart.value_axis.tick_labels.font.name = FONT
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.color.rgb = TEXT_DIM
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = LINE


def add_kpi(slide, x, y, label, value, color=ACCENT):
    box = slide.shapes.add_textbox(x, y, Inches(3.2), Inches(1.55))
    tf = box.text_frame
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = label.upper()
    r1.font.name = MONO
    r1.font.size = Pt(10)
    r1.font.color.rgb = TEXT_DIM
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.name = FONT
    r2.font.size = Pt(34)
    r2.font.bold = True
    r2.font.color.rgb = color


def slide_1_cover(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "Stillform: Applied Science", "Metacognition, EQ, composure, impulsivity control, microbias conflict")
    add_card(s, Inches(0.85), Inches(2.25), Inches(11.9), Inches(4.0))
    txt = s.shapes.add_textbox(Inches(1.2), Inches(2.75), Inches(11.2), Inches(3.2))
    tf = txt.text_frame
    points = [
        "This deck is product-mechanism-first, not generic wellness benchmarking.",
        "Each domain includes: global issue context -> Stillform mechanism -> measurable telemetry signal.",
        "No psychiatric treatment claims. Stillform is a composure training and decision-quality system.",
    ]
    for i, line in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
    add_source(s, "Stillform science sheet + app instrumentation model")


def slide_2_global_context(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "Global Composure Context", "Pressure is high; capacity is uneven")

    add_card(s, Inches(0.75), Inches(2.1), Inches(6.0), Inches(4.65))
    neg = CategoryChartData()
    neg.categories = ["Worry", "Stress", "Physical pain", "Sadness", "Anger"]
    neg.add_series("Daily negative experiences (%)", (39, 37, 32, 26, 22))
    c1 = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(1.0), Inches(2.55), Inches(5.55), Inches(3.95), neg).chart
    style_chart(c1, legend=False)
    c1.value_axis.maximum_scale = 45
    p1 = c1.series[0]
    p1.format.fill.solid()
    p1.format.fill.fore_color.rgb = RED
    p1.format.line.color.rgb = RED

    add_card(s, Inches(6.95), Inches(2.1), Inches(5.85), Inches(4.65))
    pos = CategoryChartData()
    pos.categories = ["Respect", "Laughter", "Enjoyment", "Well-rested", "Learned something"]
    pos.add_series("Daily positive experiences (%)", (88, 73, 73, 72, 52))
    c2 = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(7.2), Inches(2.55), Inches(5.4), Inches(3.95), pos).chart
    style_chart(c2, legend=False)
    c2.value_axis.maximum_scale = 100
    p2 = c2.series[0]
    p2.format.fill.solid()
    p2.format.fill.fore_color.rgb = GREEN
    p2.format.line.color.rgb = GREEN

    add_source(s, "Gallup State of the World's Emotional Health (2024)")


def slide_3_applied_science_map(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "Applied Science Map", "Domain -> global issue -> Stillform mechanism -> telemetry")

    table = s.shapes.add_table(6, 4, Inches(0.65), Inches(2.0), Inches(12.1), Inches(4.95)).table
    widths = [Inches(2.0), Inches(3.3), Inches(3.35), Inches(3.45)]
    for i, w in enumerate(widths):
        table.columns[i].width = w

    rows = [
        ("Domain", "Global issue context", "Stillform mechanism", "Measured signal"),
        ("Metacognition", "High worry/stress narrows cognitive monitoring", "Pulse labeling + Reframe pattern reflection", "Pre/post composure delta trend"),
        ("EQ", "Interpersonal reads degrade under load", "Bio-filter + interpersonal microbias prompts", "Conflict-triggered session outcomes"),
        ("Composure", "Daily regulation instability", "Morning baseline + intervention + EOD close", "14d loop completion"),
        ("Impulsivity", "Fast state -> fast action errors", "Somatic interrupt + default pathway routing", "Time-to-stabilization, post-rating gain"),
        ("Microbias conflict", "20-30% overread of others' negative intensity", "Bias detection (projection/attribution/contagion)", "Bias-tag recurrence reduction"),
    ]

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE if r == 0 else (BG if r % 2 else SURFACE_ALT)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = FONT
            run.font.size = Pt(11 if r == 0 else 10)
            run.font.bold = r == 0
            run.font.color.rgb = ACCENT if r == 0 else TEXT

    add_source(s, "Stillform science sheet; App.jsx loop, rating, and nudge telemetry")


def slide_4_metacognition(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "Metacognition", "From unobserved state loops to monitored response choice")

    add_card(s, Inches(0.75), Inches(2.1), Inches(8.1), Inches(4.65))
    d = CategoryChartData()
    d.categories = ["Worry load", "Stress load", "Learning signal", "Affect labeling support"]
    d.add_series("Issue/lever index", (39, 37, 52, 60))
    c = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.05), Inches(2.45), Inches(7.55), Inches(3.95), d).chart
    style_chart(c, legend=False)
    c.value_axis.maximum_scale = 65
    for i, p in enumerate(c.series[0].points):
        clr = [RED, RED, GREEN, BLUE][i]
        p.format.fill.solid()
        p.format.fill.fore_color.rgb = clr
        p.format.line.color.rgb = clr

    add_card(s, Inches(9.1), Inches(2.1), Inches(3.7), Inches(4.65))
    add_kpi(s, Inches(9.4), Inches(2.45), "Core mechanism", "Affect labeling", BLUE)
    add_kpi(s, Inches(9.4), Inches(3.75), "Tool", "Pulse + Reframe", GREEN)
    add_kpi(s, Inches(9.4), Inches(5.0), "Outcome", "Clearer choices", ACCENT)

    add_source(s, "Flavell 1979; Lieberman 2007; Torre & Lieberman 2018; Gallup 2024")


def slide_5_eq_microbias(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "EQ + Microbias Conflict", "When state is noisy, social interpretation gets distorted")

    add_card(s, Inches(0.75), Inches(2.1), Inches(8.1), Inches(4.65))
    d = CategoryChartData()
    d.categories = [
        "Negative intensity overread (low)",
        "Negative intensity overread (high)",
        "Daily anger",
        "Daily sadness",
    ]
    d.add_series("Observed range / prevalence (%)", (20, 30, 22, 26))
    c = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.05), Inches(2.45), Inches(7.55), Inches(3.95), d).chart
    style_chart(c, legend=False)
    c.value_axis.maximum_scale = 35
    for i, p in enumerate(c.series[0].points):
        clr = [PURPLE, PURPLE, RED, RED][i]
        p.format.fill.solid()
        p.format.fill.fore_color.rgb = clr
        p.format.line.color.rgb = clr

    add_card(s, Inches(9.1), Inches(2.1), Inches(3.7), Inches(4.65))
    txt = s.shapes.add_textbox(Inches(9.4), Inches(2.45), Inches(3.2), Inches(3.9))
    tf = txt.text_frame
    lines = [
        "Stillform applies one bias check at a time:",
        "projection, attribution, contagion, impact gap, intensity amplification.",
        "Goal: reduce reactive conflict reads and improve decision tone.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(16)
        p.font.bold = i == 0
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)

    add_source(s, "Stillform science sheet microbias section; Nature Communications 2025; Gallup 2024")


def slide_6_impulsivity(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "Impulsivity Under Load", "High arousal + low monitoring increases reactive action risk")

    add_card(s, Inches(0.75), Inches(2.1), Inches(6.0), Inches(4.65))
    d1 = CategoryChartData()
    d1.categories = ["2011 baseline", "2019 level"]
    d1.add_series("Riots/strikes/anti-gov demos index", (100, 344))
    c1 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.0), Inches(2.45), Inches(5.55), Inches(3.95), d1).chart
    style_chart(c1, legend=False)
    c1.value_axis.maximum_scale = 360
    for p in c1.series[0].points:
        p.format.fill.solid()
        p.format.fill.fore_color.rgb = RED
        p.format.line.color.rgb = RED

    add_card(s, Inches(6.95), Inches(2.1), Inches(5.85), Inches(4.65))
    d2 = CategoryChartData()
    d2.categories = ["Stress", "Worry"]
    d2.add_series("Daily pressure prevalence (%)", (37, 39))
    c2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.2), Inches(2.45), Inches(5.4), Inches(3.95), d2).chart
    style_chart(c2, legend=False)
    c2.value_axis.maximum_scale = 45
    for p in c2.series[0].points:
        p.format.fill.solid()
        p.format.fill.fore_color.rgb = ACCENT
        p.format.line.color.rgb = ACCENT

    add_source(s, "Gallup/IEP unrest trend cited in Gallup 2025 report; Gallup 2024 daily stress/worry")


def slide_7_stillform_mechanism(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "Stillform Mechanism Architecture", "How applied science is operationalized in-product")

    steps = [
        ("State capture", "Morning + bio-filter", BLUE),
        ("Interrupt", "Breathe / Body Scan / Somatic cue", GREEN),
        ("Reappraise", "AI reframe + microbias correction", ACCENT),
        ("Transfer", "One next action under control", PURPLE),
        ("Consolidate", "EOD close + pattern memory", RED),
    ]
    x0, y, w, h, gap = Inches(0.72), Inches(2.45), Inches(2.4), Inches(2.75), Inches(0.12)
    for i, (title, sub, color) in enumerate(steps):
        x = x0 + i * (w + gap)
        add_card(s, x, y, w, h, alt=i % 2 == 1)
        t = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.18), Inches(2.1), Inches(0.5))
        tp = t.text_frame.paragraphs[0]
        tr = tp.add_run()
        tr.text = title
        tr.font.name = FONT
        tr.font.size = Pt(16)
        tr.font.bold = True
        tr.font.color.rgb = color
        b = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.92), Inches(2.1), Inches(1.65))
        bp = b.text_frame.paragraphs[0]
        br = bp.add_run()
        br.text = sub
        br.font.name = FONT
        br.font.size = Pt(14)
        br.font.color.rgb = TEXT
        if i < len(steps) - 1:
            arrow = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + w - Inches(0.02), y + Inches(1.15), Inches(0.14), Inches(0.55))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEXT_DIM
            arrow.line.fill.background()

    add_source(s, "Stillform intervention flow: morning baseline -> intervention -> post-rating -> EOD loop")


def slide_8_telemetry(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "Telemetry Proof Layer", "Applied science must be visible in measurable loop outputs")

    add_card(s, Inches(0.75), Inches(2.1), Inches(8.1), Inches(4.65))
    d = CategoryChartData()
    d.categories = ["Pre/post delta", "Loop completion 14d", "Nudge recovery 14d", "Reliability score", "Tool efficacy"]
    d.add_series("System signal coverage index", (100, 100, 100, 100, 100))
    c = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.05), Inches(2.45), Inches(7.55), Inches(3.95), d).chart
    style_chart(c, legend=False)
    c.value_axis.maximum_scale = 110
    for i, p in enumerate(c.series[0].points):
        clr = [BLUE, GREEN, PURPLE, ACCENT, RED][i]
        p.format.fill.solid()
        p.format.fill.fore_color.rgb = clr
        p.format.line.color.rgb = clr

    add_card(s, Inches(9.1), Inches(2.1), Inches(3.7), Inches(4.65))
    lines = [
        "Reliability score formula:",
        "(Loop completion x 0.65) +",
        "(Nudge recovery x 0.35)",
        "",
        "Windows:",
        "14-day operating loop",
        "12-week telemetry history",
    ]
    b = s.shapes.add_textbox(Inches(9.4), Inches(2.45), Inches(3.2), Inches(3.9))
    tf = b.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(16 if i == 0 else 14)
        p.font.bold = i == 0
        p.font.color.rgb = TEXT
        p.space_after = Pt(5)

    add_source(s, "Stillform App.jsx: reliabilityScore, loop history, nudge event telemetry")


def slide_9_uat(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "UAT Decision Instrument", "Testing now scores mechanism quality, not vague sentiment")

    add_card(s, Inches(0.75), Inches(2.1), Inches(8.1), Inches(4.65))
    d = CategoryChartData()
    d.categories = [
        "Behavior changed in real moment",
        "Friction pinpoint (exact step)",
        "Reliability defect report",
        "Decision quality transfer",
        "Daily use blocker ranking",
    ]
    d.add_series("Weight (%)", (30, 20, 20, 15, 15))
    c = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.05), Inches(2.45), Inches(7.55), Inches(3.95), d).chart
    style_chart(c, legend=False)
    c.value_axis.maximum_scale = 35
    for i, p in enumerate(c.series[0].points):
        clr = [ACCENT, BLUE, GREEN, PURPLE, RED][i]
        p.format.fill.solid()
        p.format.fill.fore_color.rgb = clr
        p.format.line.color.rgb = clr

    add_card(s, Inches(9.1), Inches(2.1), Inches(3.7), Inches(4.65))
    add_kpi(s, Inches(9.4), Inches(2.45), "Questions", "5", ACCENT)
    add_kpi(s, Inches(9.4), Inches(3.75), "Mode", "Operator-grade", GREEN)
    add_kpi(s, Inches(9.4), Inches(5.0), "Purpose", "Ship decisions", BLUE)

    add_source(s, "public/uat-roadmap.html tester ask section (updated)")


def slide_10_close(prs, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_header(s, idx, total)
    add_title(s, "Closing Thesis", "Applied science -> composure skill -> better decisions")
    add_card(s, Inches(0.9), Inches(2.2), Inches(11.55), Inches(4.2))
    txt = s.shapes.add_textbox(Inches(1.35), Inches(2.7), Inches(10.8), Inches(3.35))
    tf = txt.text_frame
    lines = [
        "Stillform is not psychiatric treatment software.",
        "It is a composure and self-awareness training system.",
        "The science is applied through daily loops, state-aware interventions, and measurable behavior shift.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(29 if i == 0 else 24)
        p.font.bold = True
        p.font.color.rgb = TEXT
        p.space_after = Pt(9)
    add_source(s, "Stillform applied-science model (metacognition, EQ, composure, impulsivity, microbias)")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 10
    slide_1_cover(prs, 1, total)
    slide_2_global_context(prs, 2, total)
    slide_3_applied_science_map(prs, 3, total)
    slide_4_metacognition(prs, 4, total)
    slide_5_eq_microbias(prs, 5, total)
    slide_6_impulsivity(prs, 6, total)
    slide_7_stillform_mechanism(prs, 7, total)
    slide_8_telemetry(prs, 8, total)
    slide_9_uat(prs, 9, total)
    slide_10_close(prs, 10, total)

    out = Path(__file__).resolve().parents[1] / "public" / "Stillform_Executive_Deck.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    build()
