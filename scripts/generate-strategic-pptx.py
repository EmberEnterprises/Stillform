#!/usr/bin/env python3
"""Generate a professional, board-ready Stillform strategic deck."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BG = RGBColor(14, 15, 17)
SURFACE = RGBColor(23, 25, 29)
SURFACE_ALT = RGBColor(31, 34, 40)
LINE = RGBColor(53, 56, 64)
AMBER = RGBColor(201, 147, 58)
TEXT = RGBColor(233, 233, 234)
TEXT_DIM = RGBColor(165, 168, 176)
WHITE_SOFT = RGBColor(196, 198, 204)

TITLE_FONT = "Cormorant Garamond"
BODY_FONT = "DM Sans"
MONO_FONT = "IBM Plex Mono"


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_shell(slide, idx, total, kicker="STILLFORM · STRATEGIC DECK"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SURFACE
    bar.line.fill.background()

    k = slide.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(5.5), Inches(0.2))
    p = k.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = kicker
    r.font.name = MONO_FONT
    r.font.size = Pt(9)
    r.font.color.rgb = AMBER

    n = slide.shapes.add_textbox(Inches(11.3), Inches(0.12), Inches(1.5), Inches(0.2))
    p2 = n.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{idx:02d} / {total:02d}"
    r2.font.name = MONO_FONT
    r2.font.size = Pt(9)
    r2.font.color.rgb = TEXT_DIM


def add_title(slide, title, subtitle=None):
    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(12.0), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = TITLE_FONT
    r.font.size = Pt(40)
    r.font.color.rgb = AMBER

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.73), Inches(1.65), Inches(11.8), Inches(0.45))
        p2 = s.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = BODY_FONT
        r2.font.size = Pt(16)
        r2.font.color.rgb = TEXT_DIM


def add_section_break(slide, idx, total, section, statement):
    set_background(slide)
    add_shell(slide, idx, total)
    add_title(slide, section, "SECTION BREAK")

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(2.3), Inches(11.35), Inches(3.45))
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE
    card.line.color.rgb = LINE
    card.line.width = Pt(1.2)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(2.3), Inches(0.16), Inches(3.45))
    band.fill.solid()
    band.fill.fore_color.rgb = AMBER
    band.line.fill.background()

    b = slide.shapes.add_textbox(Inches(1.35), Inches(3.0), Inches(10.2), Inches(2.0))
    p = b.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = statement
    r.font.name = BODY_FONT
    r.font.size = Pt(26)
    r.font.color.rgb = TEXT


def add_bullets(slide, x, y, w, h, bullets, size=18, spacing=10):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(spacing)


def add_table(slide, x, y, w, h, rows, col_widths, body_size=11):
    t_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = t_shape.table
    for i, col_w in enumerate(col_widths):
        table.columns[i].width = col_w

    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = SURFACE
            else:
                cell.fill.fore_color.rgb = BG if r_idx % 2 == 1 else SURFACE_ALT
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = BODY_FONT
                    run.font.size = Pt(12 if r_idx == 0 else body_size)
                    run.font.color.rgb = AMBER if r_idx == 0 else WHITE_SOFT


def add_metric_cards(slide, cards):
    card_w = Inches(3.95)
    gap = Inches(0.25)
    x0 = Inches(0.7)
    y = Inches(2.2)
    h = Inches(3.95)
    for i, card in enumerate(cards):
        x = x0 + i * (card_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = LINE
        shape.line.width = Pt(1)

        title = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.22), card_w - Inches(0.5), Inches(0.35))
        p1 = title.text_frame.paragraphs[0]
        r1 = p1.add_run()
        r1.text = card["title"].upper()
        r1.font.name = MONO_FONT
        r1.font.size = Pt(10)
        r1.font.color.rgb = AMBER

        body = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.7), card_w - Inches(0.5), Inches(2.95))
        add_bullets(
            slide,
            x + Inches(0.25),
            y + Inches(0.7),
            card_w - Inches(0.5),
            Inches(2.95),
            card["bullets"],
            size=14,
            spacing=8,
        )


def build_deck(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 15

    # 1 cover
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 1, total)
    add_title(
        s,
        "Stillform Strategic Overview",
        "Preventive risk management for human behavior",
    )
    add_bullets(
        s,
        Inches(0.8),
        Inches(2.45),
        Inches(11.8),
        Inches(3.4),
        [
            "Mission: make composure trainable, measurable, and transferable for everyone.",
            "Category: neuroscience-guided preventive risk management for emotional maturity and decision quality.",
            "Standard: every claim maps to mechanism, evidence, and measurable product signal.",
        ],
        size=19,
    )

    # 2 section break
    s = prs.slides.add_slide(blank)
    add_section_break(s, 2, total, "Operating Thesis", "Clear category definition and operating thesis.")

    # 3 operating definition
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 3, total)
    add_title(s, "Operating definition", "Category, mechanism, and strategic position")
    add_table(
        s,
        Inches(0.7),
        Inches(2.0),
        Inches(12.0),
        Inches(4.95),
        [
            ["Dimension", "Definition", "Why it matters"],
            ["Who", "Everyone building composure as a trainable daily capability.", "Universal requirement, not niche positioning."],
            ["What", "A neuroscience-guided preventive risk management system for behavior.", "Reduces avoidable damage before it compounds."],
            ["How", "Detect state -> route intervention -> transfer action -> close daily loop -> learn patterns.", "Mechanism-driven architecture, not inspiration content."],
            ["Why", "Unregulated state degrades judgment, communication, and relationships.", "Cleaner decisions and lower downstream relational/operational cost."],
        ],
        [Inches(1.55), Inches(5.15), Inches(5.3)],
    )

    # 4 section break
    s = prs.slides.add_slide(blank)
    add_section_break(s, 4, total, "Mechanics + Science", "Science is operationalized as product behavior, not decoration.")

    # 5 mechanics map I
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 5, total)
    add_title(s, "Mechanics map I", "Science -> product implementation -> measured signal")
    add_table(
        s,
        Inches(0.55),
        Inches(1.95),
        Inches(12.25),
        Inches(5.1),
        [
            ["Mechanic", "Scientific basis", "Product implementation", "Measured signal"],
            ["Affect labeling", "Lieberman et al. (2007); Torre & Lieberman (2018)", "Pulse chips + post-session labeling", "State precision and composure-shift trend"],
            ["Cognitive reappraisal", "Ochsner & Gross (2005); Buhle et al. (2014)", "Reframe structured perspective sequence", "Pre/post composure delta"],
            ["Interoceptive regulation", "Mehling et al. (2012); Critchley & Garfinkel (2017)", "Bio-filter + body scan + somatic prompts", "Recovery latency reduction"],
            ["Implementation intentions", "Gollwitzer (1999); Gollwitzer & Sheeran (2006)", "Calibration -> default tool routing", "Action latency reduction"],
        ],
        [Inches(1.95), Inches(3.15), Inches(3.45), Inches(3.7)],
        body_size=10,
    )

    # 6 mechanics map II
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 6, total)
    add_title(s, "Mechanics map II", "Longitudinal conditioning and risk reduction")
    add_table(
        s,
        Inches(0.55),
        Inches(1.95),
        Inches(12.25),
        Inches(5.1),
        [
            ["Mechanic", "Scientific basis", "Product implementation", "Measured signal"],
            ["Stress inoculation", "Meichenbaum (1985)", "Morning baseline + EOD closure", "Loop adherence and relapse protection"],
            ["Autonomic flexibility", "Thayer & Lane (2000); Appelhans & Luecken (2006)", "Repeated regulation reps over sessions", "Faster baseline return over time"],
            ["Emotional granularity", "Barrett et al. (2001); Kashdan et al. (2015)", "Expanded emotion taxonomy in Pulse", "Higher-quality self-report and intervention fit"],
            ["Misattribution correction", "Schachter & Singer (1962); Goldstein et al. (2007)", "Bio-filter modifies interpretation before response", "False-threat interpretation reduction"],
        ],
        [Inches(1.95), Inches(3.15), Inches(3.45), Inches(3.7)],
        body_size=10,
    )

    # 7 section break
    s = prs.slides.add_slide(blank)
    add_section_break(s, 7, total, "Product System", "Preventive loop design and measurable evidence surfaces.")

    # 8 core loop cards
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 8, total)
    add_title(s, "Core loop architecture", "Six linked mechanics, one coherent daily system")
    add_metric_cards(
        s,
        [
            {
                "title": "Loop start",
                "bullets": [
                    "Morning baseline",
                    "Energy + hardware read",
                    "Risk context before action",
                    "No branch clutter",
                ],
            },
            {
                "title": "Loop middle",
                "bullets": [
                    "State-based routing",
                    "Body-first or thought-first",
                    "Reframe + regulation",
                    "Transfer to practical action",
                ],
            },
            {
                "title": "Loop close",
                "bullets": [
                    "EOD closure signal",
                    "Pattern continuity",
                    "Nudge recovery checks",
                    "Longitudinal compounding",
                ],
            },
        ],
    )

    # 9 evidence model
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 9, total)
    add_title(s, "Evidence model", "What must be true to claim impact")
    add_table(
        s,
        Inches(0.75),
        Inches(2.0),
        Inches(11.9),
        Inches(4.95),
        [
            ["Layer", "Metric", "Interpretation"],
            ["Session", "Pre/post composure delta", "Immediate regulation effect quality"],
            ["Daily loop", "Morning/EOD completion + drop-off", "Behavioral consistency under load"],
            ["Intervention", "Nudge shown/actioned/recovery", "Preventive support conversion"],
            ["Longitudinal", "Awareness latency trend", "Maturity and self-regulation speed"],
            ["Trust", "Restore reliability + entitlement truth", "Operational integrity in production"],
        ],
        [Inches(2.2), Inches(3.9), Inches(5.8)],
    )

    # 10 section break
    s = prs.slides.add_slide(blank)
    add_section_break(s, 10, total, "Business + Distribution", "Focused economics and trust-first channel strategy.")

    # 11 business model
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 11, total)
    add_title(s, "Business model", "Simple model, trust-aligned economics")
    add_table(
        s,
        Inches(0.7),
        Inches(2.0),
        Inches(12.0),
        Inches(4.9),
        [
            ["Component", "Current position", "Defensibility logic"],
            ["Revenue model", "Single tier ($14.99 / $9.99 annualized)", "Low friction + no upsell trust tax"],
            ["Category", "Preventive risk management for behavior", "Distinct from passive wellness consumption"],
            ["Distribution", "Coaches, therapist-adjacent, trusted communities", "Trust transfer outperforms cold paid acquisition"],
            ["Retention thesis", "Daily loop utility + proof of change + continuity", "Habit + evidence + context memory"],
        ],
        [Inches(2.15), Inches(4.2), Inches(5.65)],
    )

    # 12 phase cards
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 12, total)
    add_title(s, "Execution phases", "Non-calendar phase logic")
    add_metric_cards(
        s,
        [
            {
                "title": "Phase 1-2",
                "bullets": [
                    "Reliability + proof",
                    "Subscription truth",
                    "Daily loop hardening",
                    "Measured retention baseline",
                ],
            },
            {
                "title": "Phase 3",
                "bullets": [
                    "Distribution scale",
                    "Coach/practitioner channels",
                    "Referral systemization",
                    "Conversion quality gates",
                ],
            },
            {
                "title": "Phase 4",
                "bullets": [
                    "Calendar + health hookups",
                    "Consent/revoke integrity",
                    "Lower manual input burden",
                    "Context quality protection",
                ],
            },
        ],
    )

    # 13 section break
    s = prs.slides.add_slide(blank)
    add_section_break(s, 13, total, "Integrity + Risk", "Release discipline is part of the product, not an afterthought.")

    # 14 integrity + risk split
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 14, total)
    add_title(s, "Integrity operations", "How trust is enforced in code and release practice")
    left = slide_left = Inches(0.7)
    card_w = Inches(5.85)
    y = Inches(2.05)
    h = Inches(4.95)

    c1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, card_w, h)
    c1.fill.solid()
    c1.fill.fore_color.rgb = SURFACE
    c1.line.color.rgb = LINE
    c1.line.width = Pt(1)
    add_bullets(
        s,
        left + Inches(0.25),
        y + Inches(0.3),
        card_w - Inches(0.45),
        h - Inches(0.4),
        [
            "SHIP preflight gates: build + invariant + trust checks.",
            "Audience lock: universal framing; no niche-only drift.",
            "Claim lock: mechanism-first language; no overclaiming.",
            "Escalation: any invariant failure blocks release.",
            "Diagnostics on trust-critical flows (subscription/sync/restore).",
        ],
        size=14,
        spacing=8,
    )

    c2x = Inches(6.8)
    c2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c2x, y, card_w, h)
    c2.fill.solid()
    c2.fill.fore_color.rgb = SURFACE_ALT
    c2.line.color.rgb = LINE
    c2.line.width = Pt(1)
    add_bullets(
        s,
        c2x + Inches(0.25),
        y + Inches(0.3),
        card_w - Inches(0.45),
        h - Inches(0.4),
        [
            "Risk: feature sprawl -> Mitigation: ecosystem-first gate.",
            "Risk: reliability regression -> Mitigation: preflight + smoke checks.",
            "Risk: trust mismatch -> Mitigation: copy/behavior parity review.",
            "Risk: retention decay -> Mitigation: loop adherence + nudge recovery telemetry.",
            "Risk: evidence drift -> Mitigation: measured outcomes before stronger claims.",
        ],
        size=14,
        spacing=8,
    )

    # 15 close
    s = prs.slides.add_slide(blank)
    set_background(s)
    add_shell(s, 15, total)
    add_title(s, "Stillform", "Preventive risk management for human behavior")
    close = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.8), Inches(2.7))
    tf = close.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = (
        "Mission: reduce avoidable human damage by making composure trainable, measurable, and transferable.\n\n"
        "Standard: honor, honesty, and integrity in product behavior, language, and release decisions.\n\n"
        "stillformapp.com"
    )
    r.font.name = BODY_FONT
    r.font.size = Pt(24)
    r.font.color.rgb = TEXT

    prs.save(path)


if __name__ == "__main__":
    build_deck("public/Stillform_Strategic_Overview.pptx")
